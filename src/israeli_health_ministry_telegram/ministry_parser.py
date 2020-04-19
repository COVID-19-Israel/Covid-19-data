import tabula
from pptx import Presentation
import os
import pandas as pd
import re

CSV_SUFFIX = ".csv"
SPECIFIC_TABLE_PREFIX = "_table_no_"
TABLES_PARSERS_DIR = '../files_tables_parser'
DAILY_UPDATE_TEMPLATE_PATH = os.path.join(TABLES_PARSERS_DIR,'templates','daily_update_template.json')

DAILY_UPDATE_FILE_PREFIX = "מכלול_אשפוז_דיווח"
DAILY_UPDATE_TABLE_TOP_BUFFER = 2
DAILY_UPDATE_TABLE_BOTTOM_BUFFER = 0
DAILY_UPDATE_OUTPUT_DIR = "daily_update_new"
TABLE_PREFIX = "table"
TABLE_1_FULL_PREFIX = "table_1_"
TABLE_2_FULL_PREFIX = "table_2_"
PPTX_SUFFIX = "pptx"
PDF_SUFFIX = "pdf"

BLACKLIESTED_FILES = ['1200_2020-03-20','2100_2020-03-29',
                      '1100_2020-03-21','2000_2020-03-28']
TABLE_0_COLUMNS = ['נפטר','קשה כעת','בינוני']
RECOVERED_TEXT = 'החלימו ושוחררו'
PARTS_OF_CONFIRMED = ['Mild','Moderate','Deceased','Critical','Recovered']
FIELDS_RENAMING = {
    TABLE_1_FULL_PREFIX + 'קל': 'Mild',
    TABLE_1_FULL_PREFIX + 'בינוני': 'Moderate',
    TABLE_1_FULL_PREFIX + 'נפטר': 'Deceased',
    TABLE_1_FULL_PREFIX + 'קשה כעת': 'Critical',
    TABLE_1_FULL_PREFIX + 'קשה  כעת': 'Critical',
    TABLE_1_FULL_PREFIX + 'סה"כ מחלימים': 'Recovered',
    TABLE_1_FULL_PREFIX + 'ה"כ מחלימים': 'Recovered',
    TABLE_2_FULL_PREFIX + 'החלימו ושוחררו': 'Recovered',
    TABLE_2_FULL_PREFIX + 'מחלימים': 'Recovered',
    TABLE_2_FULL_PREFIX + 'Unnamed: 0': 'Recovered',
}
def format_int(value: str) -> int:
    try:
        return(int(float(value)))
    except:
        if value in ['','nan', None]:
            return None
        else:
            value = value.strip().replace(',', '').replace('+', '')
            return None if value == '' else int(float(value))

def get_time_from_filename(file_name):
    HOUR_STRING = "לשעה"
    path_parts = re.split("[_.]", file_name)
    hour = path_parts[path_parts.index(HOUR_STRING) + 1]
    date = path_parts[0]
    time = '_'.join([hour, date])
    return time

def connect_csvs(file_1, file1_prefix, file_2, file2_prefix):
    df_1_headers = [file1_prefix + val for val in file_1.columns]
    df_2_headers = [file2_prefix + val for val in file_2.columns]
    concat_table = pd.concat([file_1, file_2], axis=1)
    concat_table.columns = df_1_headers + df_2_headers
    return concat_table

class MinistryFileParser:
    def __init__(self, path, output_dir):
        self.output_dir = os.path.join(output_dir,DAILY_UPDATE_OUTPUT_DIR)
        self.path = path
        self._data = None

    def run(self):
        file_name = os.path.basename(self.path).split(".")
        file_suffix = file_name[-1]
        file_name = "".join(file_name[:-1])
        self.file_name = file_name
        if DAILY_UPDATE_FILE_PREFIX in file_name:
            if(get_time_from_filename(file_name)) in BLACKLIESTED_FILES:
                return
            if(os.path.exists(self._create_output_file_path())):
                return
            try:
                if PPTX_SUFFIX == file_suffix:
                    self.parse_pptx_file()
                elif PDF_SUFFIX == file_suffix:
                    self.parse_pdf_file()

                self.data_reformatting()

                self.export_to_csv()
            except:
                print("could not parse {}".format(get_time_from_filename(file_name)))

    def data_reformatting(self):
        self._data.columns = [FIELDS_RENAMING.get(val,None) for val in self._data.columns]

        while None in self._data.columns:
            self._data = self._data[self._data.columns.dropna()]

        self._data = self._data.applymap(format_int)
        confirmed = sum([int(self._data.get(val,[0])[0]) for val in PARTS_OF_CONFIRMED])
        self._data['Confirmed'] = [confirmed]

    def parse_pptx_file(self):
        prs = Presentation(self.path)
        prs_tables = list()

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                parsed_table = list()
                tbl = shape.table
                for row_index in range(0, len(tbl.rows)):
                    tbl_row = list()
                    for col_index in range(0, len(tbl.columns)):
                        table_cell = tbl.cell(row_index, col_index)
                        cell_data = self._extract_data_from_cell(table_cell)
                        tbl_row.append(cell_data)
                    parsed_table.append(tbl_row)
                prs_tables.append(parsed_table)

        prs_tables = prs_tables[DAILY_UPDATE_TABLE_BOTTOM_BUFFER:DAILY_UPDATE_TABLE_TOP_BUFFER]
        table_values = list()
        table_keys = list()

        for idx,table in enumerate(prs_tables):
            for row_index in range(len(table)):
                for col_index in range(len(table[0])):
                    if len(table_keys) >=14:
                        break
                    if table[row_index][col_index].replace(",", "").isdigit():
                        table_values.append(table[row_index][col_index])
                        key = self._find_key_by_value(table, row_index, col_index)
                        table_number = '1' if len(table_keys)<6 else '2'
                        key = "_".join([TABLE_PREFIX,table_number,key])
                        table_keys.append(key)
        self._data = pd.DataFrame(columns=table_keys, data=[table_values])

    @staticmethod
    def _extract_data_from_cell(table_cell):
        data = list()
        cell_paragraphs = table_cell.text_frame.paragraphs
        for cell_paragraph in cell_paragraphs:
            for run in cell_paragraph.runs:
                data.append(run.text)
        return " ".join(data)

    @staticmethod
    def _find_key_by_value(table, row_index, col_index):
        if (
                table[row_index - 1][col_index]
                and not table[row_index - 1][col_index].replace(",", "").isdigit()
        ):
            return table[row_index - 1][col_index]
        if (
                table[row_index][col_index - 1]
                and not table[row_index][col_index - 1].replace(",", "").isdigit()
        ):
            return table[row_index][col_index - 1]


    def parse_pdf_file(self):
        pdf_tables = tabula.read_pdf_with_template(
            input_path=self.path, template_path=DAILY_UPDATE_TEMPLATE_PATH)
        # for table in pdf_tables:
        #     print (table)
        #     print ("@@@@@@")

        pdf_table0 = self._fix_critical_confirmed_table(pdf_tables[0])

        #print(pdf_table0)
        table_no_1 = pd.concat([pdf_tables[2], pdf_table0, pdf_tables[1]], axis=1)
        table_no_1 = self.table_1_hotfixes(table_no_1)
        if 'ה"כ מחלימים' in table_no_1.columns or 'סה"כ מחלימים' in table_no_1.columns:
            table_no_1.columns = ["table_1_" + val for val in table_no_1.columns]
            self._data = table_no_1
        else:
            table_no_2 = pd.concat([pdf_tables[5], pdf_tables[4]], axis=1)
            table_no_2 = self.table_2_hotfixes(table_no_2)
            concat_table = connect_csvs(table_no_1, "table_1_", table_no_2, "table_2_")
            self._data = concat_table

    @staticmethod
    def _fix_critical_confirmed_table(critical_confirmed_table):
        try:
            fixed_critical_confirmed_table = critical_confirmed_table.drop(
        critical_confirmed_table.index[0])
        except:
            fixed_critical_confirmed_table = critical_confirmed_table
        if len(fixed_critical_confirmed_table.columns)>0:
            if fixed_critical_confirmed_table.columns[0].isdigit():
                return pd.DataFrame(columns=TABLE_0_COLUMNS,
                                    data=[list(fixed_critical_confirmed_table.columns)])

        return fixed_critical_confirmed_table.reset_index(drop=True)

    def table_1_hotfixes(self,table):
        if get_time_from_filename(self.file_name) == '2100_2020-03-23':
            table['קשה כעת'] = [re.sub("[^0-9]", "",table['קשה כעת'][0])]
        return table

    def table_2_hotfixes(self,table):
        if get_time_from_filename(self.file_name) in \
                ['2100_2020-03-24','0800_2020-03-23','2100_2020-03-23','0800_2020-03-25']:
            return pd.DataFrame(columns=[RECOVERED_TEXT],data=[table.columns[1]])
        return table

    def _create_output_file_path(self):
        file_name = os.path.basename(self.path)
        file_name = "".join(file_name.split(".")[:-1])
        file_name = get_time_from_filename(file_name)
        return os.path.join(self.output_dir,file_name)+CSV_SUFFIX

    def export_to_csv(self):
        os.makedirs(self.output_dir, exist_ok=True)
        output_file_name = self._create_output_file_path()
        self._data.to_csv(output_file_name, index=False, encoding="utf-8")
