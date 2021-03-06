"""
Purpose: This module represents file parsers that parses the files' tables into csv format file.
"""

import tabula
from pptx import Presentation
import os
import pandas as pd
import json
import logging
from parser_translator import ParserTranslator
from logger import create_log

FIELD_SEP = "@@@"
CSV_SUFFIX = ".csv"
SPECIFIC_TABLE_PREFIX = "_table_no_"

CURRENT_DIR = os.path.dirname(__file__)

DAILY_UPDATE_TEMPLATE_PATH = os.path.join(CURRENT_DIR, "templates\\daily_update_template.json")
DENMARK_FILES_TEMPLATE_PATH = os.path.join(CURRENT_DIR, 'templates\\denmark_files_template.json')

OLD_CITIES_COLUMNS = {"ישוב", "אוכלוסיה נכון ל 2018-", "מספר חולים"}
CITIES_HEADER_KEYWORDS = {"אוכלוסיה", "חולים", "מאומתים", "עיר", "מחלימים"}
DAILY_UPDATE_FILE_PREFIX = "מכלול_אשפוז_דיווח"
DENMARK_FILE_PREFIX = 'covid19-overvaagningsrapport-'

DOWNLOADED_FILES_DICT_PATH = os.path.join(CURRENT_DIR, r"..\query_script\data\MOHreport_DOWNLOADED.json")
FILES_BLACKLIST_PATH = os.path.join(CURRENT_DIR, r"files_blacklist.txt")

DAILY_UPDATE_TABLE_TOP_BUFFER = 1
DAILY_UPDATE_TABLE_BOTTOM_BUFFER = 0


class FileParser:
    """
    This class represents a General File Parser
    """

    PPTX_SUFFIX = "pptx"
    PDF_SUFFIX = "pdf"
    XLSX_SUFFIX = "xlsx"
    
    def __init__(self, path, output_dir=str()):
        """
        Saves 3 parameters about a the file:
        path- The file path
        _data- The parsed data. matrixs or DataFrame
        output_dir - Dir to save the csv's in.
        :param path: the file path to parse
        """
        self.output_dir = output_dir
        self.path = path
        self._data = list()

    def run(self):
        """
        This function parses the input file, and exports it to csv
        :param export_file_to_csv - is file should be exported to csv file
        :return: None
        """
        file_name = os.path.basename(self.path).split(".")
        file_suffix = file_name[-1]
        file_name = "".join(file_name[:-1])

        with open(FILES_BLACKLIST_PATH, mode="r+") as file:
            blacklist_files = file.read().split("\n")

        exists_csv_files = [
            f for dp, dn, filenames in os.walk(self.output_dir) for f in filenames
        ]

        # Skips on the file, if there is an output file for it already
        if [
            file_name
            for csv_file in exists_csv_files
            if csv_file.startswith(file_name + SPECIFIC_TABLE_PREFIX)
        ]:
            logging.warning(f"file: {file_name} has already been parsed")
            return

        if os.path.basename(self.path) in blacklist_files:
            logging.warning(f"{os.path.basename(self.path)} is in the files' blacklist")
            return

        # Matches correct parser by file type.
        if FileParser.PPTX_SUFFIX == file_suffix:
            parser = PptxParser(self.path, self.output_dir)
        elif FileParser.PDF_SUFFIX == file_suffix:
            parser = PdfParser(self.path, self.output_dir)
        elif FileParser.XLSX_SUFFIX == file_suffix:
            logging.warning("We don't have parser for xlsx files yet")
            return
        else:
            raise ValueError(
                f"This class don't parse files of the type: {file_suffix}, "
                f"path of the file: {os.path.basename(self.path)}"
            )

        parser.parse_file()
        parser.export_to_csv()

        return parser._data

    def parse_file(self):
        """
        This function parses a pptx file into list of tables (table=matrix)
        :return: None
        """
        raise NotImplementedError

    def _create_output_file_path(self, table_index):
        """
        This function will create the output file, if doesn't exist, and return it's path
        :param table_index: the index of the table in the input file
        :return: the relative path of the output file
        """
        file_name = os.path.basename(self.path)
        file_name = "".join(file_name.split(".")[:-1])
        output_file_name = "".join(
            [
                self.output_dir,
                file_name,
                SPECIFIC_TABLE_PREFIX,
                str(table_index),
                CSV_SUFFIX,
            ]
        )

        with open(output_file_name, mode="w+"):
            pass
        return output_file_name

    def export_to_csv(self):
        """
        gets a list of tables (table = matrix)
        exports each one to different csv
        :return: None
        """
        os.makedirs(self.output_dir, exist_ok=True)

        if not self._data or len(self._data) == 0:
            logging.info("Didn't parse any table from this file.")
            try:
                with open(FILES_BLACKLIST_PATH, mode="a") as file:
                    file.write(os.path.basename(self.path) + "\n")
            except Exception:
                logging.error(
                    f"Failed to add {os.path.basename(self.path)} into blacklist- unreadable chars."
                )
            return

        logging.info(f"Got {len(self._data)} tables to export")
        for table_index, table in enumerate(self._data, start=1):
            if type(table) == list:
                table_df = pd.DataFrame(columns=table[0], data=table[1:])
            # type is DataFrame
            else:
                table_df = table
            output_file_name = self._create_output_file_path(table_index)
            table_df.to_csv(output_file_name, index=False, encoding="utf-8")
            logging.info(f"Exported: {output_file_name} .")


class PptxParser(FileParser):
    """
    This class represents a pptx file parser.
    """

    DAILY_UPDATE_OUTPUT_DIR = "daily_update\\"

    def parse_file(self):
        prs = Presentation(self.path)
        prs_tables = PptxParser._parse_tables_from_pres(prs)

        prs_tables = self._parse_daily_update(prs_tables)
        self._data = prs_tables

    @staticmethod
    def _parse_tables_from_pres(prs):
        """
        Extracts all tables in the presentation to a list of matrices.
        :param prs: Presentation Object
        :return: list of matrices
        """
        prs_tables = list()

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                parsed_table = list()
                tbl = shape.table
                for row_index in range(0, len(tbl.rows)):
                    tbl_row = list()
                    for col_index in range(0, len(tbl.columns)):
                        table_cell = tbl.cell(row_index, col_index)
                        cell_data = PptxParser._extract_data_from_cell(table_cell)
                        tbl_row.append(cell_data)
                    parsed_table.append(tbl_row)
                prs_tables.append(parsed_table)

        logging.info(f"Parsed {len(prs_tables)} tables from presentation")
        return prs_tables

    @staticmethod
    def _extract_data_from_cell(table_cell):
        """
        This function extracts the data from pptx table's cell
        :param table_cell: the pptx table's cell
        :return: the data
        """
        translator = ParserTranslator(to_lang="en", from_lang="he")
        data = list()
        cell_paragraphs = table_cell.text_frame.paragraphs
        for cell_paragraph in cell_paragraphs:
            for run in cell_paragraph.runs:
                data.append(run.text)
        return translator.translate_word(" ".join(data))

    def _parse_daily_update(self, tables):
        """
        This function Checks if the pptx file is a file contains COVID-19 daily update.
        :return: None
        """
        parsed_tables = tables
        if DAILY_UPDATE_FILE_PREFIX in os.path.basename(self.path):
            logging.info("Detected Daily Update PPTX structure.")
            parsed_tables = DailyUpdatePptxParser.parse_file(tables)
            self.output_dir = os.path.join(self.output_dir, PptxParser.DAILY_UPDATE_OUTPUT_DIR)
            logging.info("Finished Daily Update PPTX parse.")
        return parsed_tables


class DailyUpdatePptxParser(PptxParser):
    """
    Parser for daily Update in PPTX format.
    """

    @staticmethod
    def parse_file(tables):
        """
        Fix data to be in specific structure.
        :param tables:
        :return:
        """
        parsed_tables = list()
        for table in tables[
            DAILY_UPDATE_TABLE_BOTTOM_BUFFER:DAILY_UPDATE_TABLE_TOP_BUFFER
        ]:
            table_values = list()
            table_keys = list()
            for row_index in range(len(table)):
                for col_index in range(len(table[0])):
                    if table[row_index][col_index].replace(",", "").isdigit():
                        table_values.append(table[row_index][col_index])
                        table_keys.append(
                            DailyUpdatePptxParser._find_key_by_value(
                                table, row_index, col_index
                            )
                        )
            parsed_tables.append([table_keys, table_values])
        return parsed_tables

    @staticmethod
    def _find_key_by_value(table, row_index, col_index):
        """
        This function finds the title of the value (number) inside the table
        :param table - the table
        :param row_index - the index of the row
        :param col_index - the index of the column
        :return: the title of the value
        """
        try:
            # checks if there is a title on top of the value
            if (
                table[row_index - 1][col_index]
                and not table[row_index - 1][col_index].replace(",", "").isdigit()
            ):
                return table[row_index - 1][col_index]
            # checks if there is a title on right of the value
            if (
                table[row_index][col_index - 1]
                and not table[row_index][col_index - 1].replace(",", "").isdigit()
            ):
                return table[row_index][col_index - 1]
            raise ValueError(
                f"You have in the cell: {row_index},{col_index} a number with"
                f"the value: {table[row_index][col_index]} without a title"
                f"(titles are supposed to be on top or on the right of the number)"
            )
        except IndexError:
            raise ValueError(
                f"You have in the cell: {row_index},{col_index} a number with"
                f"the value: {table[row_index][col_index]} without a title"
                f"(titles are supposed to be on top or on the right of the number)"
            )


class PdfParser(FileParser):
    """
    This class represents a parser for every Pdf file.
    """

    CITIES_OUTPUT_DIR = "cities\\"
    DAILY_UPDATE_OUTPUT_DIR = "daily_update\\"
    DENMARK_OUTPUT_DIR = 'denmark_data\\'

    def parse_file(self):
        try:
            pdf_tables = tabula.read_pdf(
                input_path=self.path, pages="all", stream=True, silent=True
            )
        except Exception:
            logging.error(f"failed to read {os.path.basename(self.path)}")
            return
        pdf_parse_functions = [self._parse_old_cities,
                               self._parse_daily_update,
                               self._parse_cities,
                               self._parse_denmark_data]

        parse_successed = False
        for parse_function in pdf_parse_functions:
            if parse_successed:
                break
            parse_successed = parse_function()

    def _parse_old_cities(self):
        """
        table type check: checks if any of the headers match the cities format.
        if yes, parses all.
        :return: is function succeeded
        """
        pdf_tables = tabula.read_pdf(
            input_path=self.path, pages="all", stream=True, silent=True
        )
        for table_df in pdf_tables:
            if OLD_CITIES_COLUMNS.issubset(set(table_df.columns.tolist())):
                logging.info("Detected Cities PDF structure.")
                parser = CitiesOldPdfParser(self.path)
                logging.info("Finished Cities PDF parse.")
                self._data.append(parser.parse_file())
                self.output_dir = os.path.join(self.output_dir, PdfParser.CITIES_OUTPUT_DIR)
                return True
        return False

    def _parse_daily_update(self):
        """
        This function Checks if the pdf file is a file contains COVID-19 daily update.
        :return: is function succeeded
        """
        if DAILY_UPDATE_FILE_PREFIX in os.path.basename(self.path):
            logging.info("Detected Daily Update PDF structure.")
            parser = DailyUpdatePdfParser(self.path)
            self._data = parser.parse_file()
            logging.info("Finished Daily Update PDF parse.")
            self.output_dir = os.path.join(self.output_dir, PdfParser.DAILY_UPDATE_OUTPUT_DIR)
            return True
        return False

    def _parse_cities(self):
        """
        This function checks if the pdf file is a file contains COVID-19 data divided by israeli cities.
        :return: is function succeeded
        """
        pdf_tables = tabula.read_pdf(
            input_path=self.path, pages="all", stream=True, silent=True
        )
        if not pdf_tables or 3 > len(pdf_tables[0]):
            return False
        first_four_lines = [pdf_tables[0].columns.tolist()] + pdf_tables[0][:5].values.tolist()
        header_words = {val for line in first_four_lines for val in line}
        if CITIES_HEADER_KEYWORDS.issubset(header_words):
            logging.info("Detected new Cities PDF structure.")
            parser = CitiesPdfParser(self.path)
            logging.info("Finished new Cities PDF parse.")
            self._data.append(parser.parse_file())
            self.output_dir = os.path.join(self.output_dir, PdfParser.CITIES_OUTPUT_DIR)
            return True
        return False

    def _parse_denmark_data(self):
        """
        This function checks if the pdf file is a file contains denmark's COVID-19 daily update data.
        :return: is function succeeded
        """
        if DENMARK_FILE_PREFIX in self.path:
            logging.info('Detected Denmark Update PDF structure.')
            parser = DenmarkPdfParser(self.path)
            self._data = parser.parse_file()
            logging.info('Finished Denmark Update PDF parse.')
            self._output_dir = os.path.join(self.output_dir, PdfParser.DENMARK_OUTPUT_DIR)
            return True
        return False

    @staticmethod
    def _are_rows_completed(matrix, row_index, previous_row_offset, values_end_col):
        """
        This function checks if the following rows can complete each other.
        If they are, it returns the merged row, else it returns an empty list
        :param matrix - the table that contains the rows.
        :param row_index - the index of the first row.
        :param pervious_row_offset - the offset between the indexes of the rows.
        :param values_end_col - the column that the values ends at (if there is any header column in the table).
        :return: merged row if the rows can be merged, else an enpty list.
        """
        merged_row = list()
        if 0 > previous_row_offset:
            min_index = row_index + previous_row_offset
            max_index = row_index
        else:
            min_index = row_index
            max_index = row_index + previous_row_offset
        merged_values = zip(matrix[row_index][:values_end_col],
                            matrix[row_index + previous_row_offset][:values_end_col])
        for values in merged_values:
            if None not in values:
                return list()
            merged_row.append(values[1 - values.index(None)])

        for headers in zip(matrix[min_index][values_end_col:], matrix[max_index][values_end_col:]):
            merged_row.append((' '.join([str(headers[0]), str(headers[1])])
                               .replace('None ', '')).replace(' None', ''))

        return merged_row

    @staticmethod
    def _merge_completed_lines(merged_table, top_to_bottom, is_col_header):
        """
        This function merges all of the following rows that can complete each other in a table.
        :param merged_table - the table that its rows need to be merged.
        :param top_to_bottom - is the direction of the table is from top to bottom.
        :param is_col_header - is there a header column in the table.
        :return: None.
        """
        merged_row = list()
        start_index = len(merged_table) - 2
        end_index = 0
        index_jumps = -1

        if merged_table:
            values_end_col = len(merged_table[0])

        if is_col_header:
            values_end_col -= 1

        if top_to_bottom:
            start_index = 1
            end_index = len(merged_table) - 1
            index_jumps = 1

        row_index = start_index

        for i in range(start_index, end_index, index_jumps):
            merged_row = PdfParser._are_rows_completed(merged_table, row_index,(-1 * index_jumps), values_end_col)
            if merged_row:
                merged_table[row_index - index_jumps] = merged_row
                merged_table.remove(merged_table[row_index])
                if 0 > index_jumps:
                    row_index += index_jumps
            else:
                row_index += index_jumps

    @staticmethod
    def _concat_empty_lines(concated_table, is_col_header, top_to_bottom=True):
        """
        This function concats empty lines that made because of line-break in
        the table to the line before (or after, in case of empty line in the index 0)
        :param concated_table - the table that need to be concated
        :return: None
        """
        PdfParser._merge_completed_lines(concated_table, top_to_bottom, is_col_header)

        row_index = 1
        for i in range(1, len(concated_table)):
            try:
                if (None in concated_table[row_index]
                    or 'None' in concated_table[row_index]
                    or (None in concated_table[0] and 1 == row_index)):
                    full_fields = zip(concated_table[row_index - 1], concated_table[row_index])
                    for col_index, full_field in enumerate(full_fields):
                        concated_table[row_index - 1][col_index] = (' '.join([str(full_field[0]),
                                                                              str(full_field[1])])
                                                                    .replace('None ', '')).replace(' None', '')
                    concated_table.remove(concated_table[row_index])
                else:
                    row_index += 1
            except Exception:
                print(f"i: {i}, row_index: {row_index}, table: {concated_table}")

    @staticmethod
    def _translate_table(translated_table, to_lang="en", from_lang="he"):
        """
        This function translates a table from danish to english (only the top line and the bottom line
        :param translated_table - the table that need to be translated
        :return: the translated table
        """
        translator = ParserTranslator(to_lang=to_lang, from_lang=from_lang)
        for row_index in range(0, len(translated_table)):
            for col_index in range(len(translated_table[0])):
                translated_table[row_index][col_index] = translator.translate_word(
                    str(translated_table[row_index][col_index]))
        return translated_table


class CitiesOldPdfParser(PdfParser):
    """
    Parser class of a pdf file that contains COVID-19 data divided into cities.
    """

    def parse_file(self):
        translator = ParserTranslator(to_lang="en", from_lang="he")
        pdf_tables = tabula.read_pdf(
            input_path=self.path, pages="all", stream=True, silent=True
        )
        logging.info(f"Parsed {len(pdf_tables)} tables in cities PDF")

        fixed_data = []
        headers = ["City_Name", "Population", "Infected"]
        fixed_data.append(headers)

        for data_df in pdf_tables:
            data_df = data_df.where(pd.notnull(data_df), None)
            list_data = data_df.values.tolist()

            for line in list_data:
                # Solves the problem that the data moved one column right in the middle of the file
                if line[1] is not None:
                    if line[1].replace(",", "").isdigit():
                        fixed_data.append(
                            [
                                translator.translate_word(line[0]),
                                line[1].replace(",", ""),
                                str(line[2]).replace(",", ""),
                            ]
                        )
                else:
                    fixed_data.append(
                        [
                            translator.translate_word(line[0]),
                            line[2].replace(",", ""),
                            str(line[3]).replace(",", ""),
                        ]
                    )
        return fixed_data


class CitiesPdfParser(PdfParser):
    def parse_file(self):
        pdf_tables = tabula.read_pdf(input_path=self.path,
                                     pages="all",
                                     stream=True,
                                     silent=True)
        is_first_iteration = True
        top_to_bottom = True
        parsed_table = list()
        first_headers_len = len(pdf_tables[0].columns)
        for pdf_table in pdf_tables:
            if len(pdf_table.columns) != first_headers_len:
                # dont add this table.
                continue
            concated_table = list()
            pdf_table = pdf_table.where(pd.notnull(pdf_table), None)
            table_headers = [header if 'Unnamed' not in str(header)
                             else None
                             for header in pdf_table.keys()]
            if list(range(len(table_headers))) != table_headers:
                concated_table.append(table_headers)
            [concated_table.append(row) for row in pdf_table.values.tolist()]
            if is_first_iteration:
                top_to_bottom = not (None in concated_table[0] and any(concated_table[0]))  # row is not empty
                is_first_iteration = False
            PdfParser._concat_empty_lines(concated_table, is_col_header=True, top_to_bottom=top_to_bottom)

            for row in PdfParser._translate_table(concated_table):
                parsed_table.append(row)

        return parsed_table


class DailyUpdatePdfParser(PdfParser):
    """
    This class represents a parser of a pdf file that contains COVID-19 daily update data.
    """

    def parse_file(self):
        fixed_pdf_tables = list()
        translated_pdf_tables = list()
        pdf_tables = tabula.read_pdf_with_template(
            input_path=self.path, template_path=DAILY_UPDATE_TEMPLATE_PATH
        )
        logging.info(f"Parsed {len(pdf_tables)} tables in DailyUpdate PDF")
        # first 3 tables of pdf_tables are from confirmed patients table (top)
        # last 3 tables of pdf_tables are from treatment table (bottom)
        try:
            pdf_table3 = DailyUpdatePdfParser._fix_treatment_table(pdf_tables[3])
            pdf_table0 = DailyUpdatePdfParser._fix_critical_confirmed_table(
                pdf_tables[0]
            )
        except Exception:

            logging.error(f"failed to parse file: {os.path.basename(self.path)}")
            return

        # concatenating to form two tables - according to the pdf's format
        fixed_pdf_tables.append(
            pd.concat([pdf_tables[2], pdf_table0, pdf_tables[1]], axis=1)
        )
        fixed_pdf_tables.append(
            pd.concat([pdf_tables[5], pdf_table3, pdf_tables[4]], axis=1)
        )

        translator = ParserTranslator(to_lang="en", from_lang="he")
        for pdf_table in fixed_pdf_tables:
            temp_df = pd.DataFrame(
                data=[pdf_table.columns.tolist()] + pdf_table.values.tolist()
            )
            temp_df = temp_df.applymap(lambda x: translator.translate_word(str(x)))
            translated_df = pd.DataFrame(
                columns=temp_df.values.tolist()[0], data=[temp_df.values.tolist()[1]]
            )
            translated_pdf_tables.append(translated_df)

        return translated_pdf_tables

    @staticmethod
    def _fix_critical_confirmed_table(critical_confirmed_table):
        fixed_critical_confirmed_table = critical_confirmed_table.drop(
            critical_confirmed_table.index[0]
        )
        return fixed_critical_confirmed_table.reset_index(drop=True)

    @staticmethod
    def _fix_treatment_table(treatment_table):
        fixed_treatment_table = pd.DataFrame(
            columns=["number", "state"],
            data=treatment_table.values.tolist() + [treatment_table.columns.tolist()],
        )

        fixed_treatment_table = fixed_treatment_table.transpose().values.tolist()

        return pd.DataFrame(
            columns=fixed_treatment_table[1], data=[fixed_treatment_table[0]]
        )

class DenmarkPdfParser(PdfParser):
    """
    This class represents a Pdf file that contains the daily update of denmark
    """
    def parse_file(self):
        pdf_tables = tabula.read_pdf_with_template(
            input_path=self.path,
            template_path=DENMARK_FILES_TEMPLATE_PATH,
            pandas_options=dict(header=None)
        )
        parsed_tables = list()

        for pdf_table in pdf_tables:
            concated_table = list()
            pdf_table = pdf_table.where(pd.notnull(pdf_table), None)
            table_headers = [header if 'Unnamed' not in str(header) and 'Antal' != str(header)
                             else None
                             for header in pdf_table.keys()]
            if list(range(len(table_headers))) != table_headers:
                concated_table.append(table_headers)
            for row in pdf_table.values.tolist():
                concated_row = list()
                for col in row:
                    # Converts the dot to comma (in danish, they are reversed)
                    if float == type(col):
                        concated_row.append(','.join([(str(col)).split('.')[0],
                                                      ((str(col) + '000').split('.')[1][:3])]))
                    elif 'Antal' != str(col) and 'None' != str(col):
                        concated_row.append(str(col).replace('.', ','))
                    else:
                        concated_row.append(None)
                concated_table.append(concated_row)
            PdfParser._concat_empty_lines(concated_table, is_col_header=False, top_to_bottom=True)
            parsed_table = PdfParser._translate_table(concated_table)
            parsed_tables.append(parsed_table)
        return parsed_tables