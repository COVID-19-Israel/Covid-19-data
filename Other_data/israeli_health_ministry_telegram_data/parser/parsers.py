"""
Purpose: This module represents file parsers that parses the files' tables into csv format file.
"""

import tabula
from pptx import Presentation
import os
import pandas as pd
import json
from parser_translator import ParserTranslator

FIELD_SEP = '@@@'
CSV_SUFFIX = '.csv'
SPECIFIC_TABLE_PREFIX = '_table_no_'
OUTPUT_DIR = '..\\csv_files\\'
CITIES_OUTPUT_DIR = OUTPUT_DIR + 'cities\\'
DAILY_UPDATE_OUTPUT_DIR = OUTPUT_DIR + 'daily_update\\'
DAILY_UPDATE_TEMPLATE_PATH = 'templates\\daily_update_template.json'

CITIES_COLUMNS = {'ישוב', 'אוכלוסיה נכון ל 2018-', 'מספר חולים'}
DAILY_UPDATE_FILE_PREFIX = 'מכלול_אשפוז_דיווח'

DOWNLOADED_FILES_DICT_PATH = r"..\query_script\data\MOHreport_DOWNLOADED.json"

DAILY_UPDATE_TABLE_TOP_BUFFER = 2
DAILY_UPDATE_TABLE_BOTTOM_BUFFER = 0


class FileParser:
    """
    This class represents a General File Parser
    """
    PPTX_SUFFIX = 'pptx'
    PDF_SUFFIX = 'pdf'

    def __init__(self, path):
        self.path = path
        self._data = list()
        self._output_dir = str()

    def run(self):
        """
        This function parses the input file, and exports it to csv
        :return: None
        """
        file_name = os.path.basename(self.path).split('.')
        file_suffix = file_name[-1]
        file_name = ''.join(file_name[:-1])

        csv_files = [os.path.join(dp, f) for dp, dn, filenames in os.walk(OUTPUT_DIR)
                     for f in filenames]

        # Skips on the file, if there is an output file for it already
        if [file_name for csv_file in csv_files if file_name in csv_file]:
            return

        if FileParser.PPTX_SUFFIX == file_suffix:
            parser = PptxParser(self.path)
        elif FileParser.PDF_SUFFIX == file_suffix:
            parser = PdfParser(self.path)
        else:
            raise ValueError(f"This class don't parse files of the type: {file_suffix}, "
                             f"path of the file: {os.path.basename(self.path)}")

        parser.parse_file()
        parser.export_to_csv()

    def parse_file(self):
        """
        This function parses a pptx file into list of tables (table=matrix)
        :return: None
        """
        raise NotImplementedError

    def _create_output_file_path(self, table_index):
        """
        This function will create the output file, if doesn't exist, and return it's path
        :param table_index: the index of the table in the input file
        :return: the path of the output file
        """
        file_name = os.path.basename(self.path)
        file_name = "".join(file_name.split(".")[:-1])
        output_file_name = ''.join([self._output_dir,
                                    file_name,
                                    SPECIFIC_TABLE_PREFIX,
                                    str(table_index),
                                    '_',
                                    self._get_file_date(),
                                    CSV_SUFFIX])

        print(f"exported: {output_file_name}")
        with open(output_file_name, mode='w+'):
            pass
        return output_file_name

    def export_to_csv(self):
        """
        gets a list of tables (table = matrix)
        exports each one to different csv
        :return: None
        """
        for table_index, table in enumerate(self._data, start=1):
            if type(table) == list:
                table_df = pd.DataFrame(columns=table[0], data=table[1:])
            # type is DataFrame
            else:
                table_df = table
            output_file_name = self._create_output_file_path(table_index)
            table_df.to_csv(output_file_name, index=False, encoding='utf-8')

    def _get_file_date(self):
        filename = os.path.basename(self.path)
        with open(DOWNLOADED_FILES_DICT_PATH, "r") as f:
            downloaded_files_dict = json.load(f)
        return downloaded_files_dict[filename]


class PptxParser(FileParser):
    """
    This class represents a pptx file parser.
    """
    @staticmethod
    def _extract_data_from_cell(table_cell):
        """
        This function extracts the data from pptx table's cell
        :param table_cell : the pptx table's cell
        :return: the data
        """
        translator = ParserTranslator(to_lang='en', from_lang='he')
        data = list()
        cell_paragraphs = table_cell.text_frame.paragraphs
        for cell_paragraph in cell_paragraphs:
            for run in cell_paragraph.runs:
                data.append(run.text)
        return translator.translate_word(' '.join(data))

    def _parse_daily_update(self, tables):
        """
        This function Checks if the pptx file is a file contains COVID-19 daily update.
        :return: None
        """
        parsed_tables = tables
        if os.path.basename(self.path).startswith(DAILY_UPDATE_FILE_PREFIX):
            print("daily update parse")
            parsed_tables = DailyUpdatePptxParser.parse_file(tables)
            self._output_dir = DAILY_UPDATE_OUTPUT_DIR
        return parsed_tables

    def parse_file(self):
        prs = Presentation(self.path)
        prs_tables = list()

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_table:
                    continue
                parsed_table = list()
                tbl = shape.table
                for row_index in range(0, len(tbl.rows)):
                    tbl_row = list()
                    for col_index in range(0, len(tbl.columns)):
                        table_cell = tbl.cell(row_index, col_index)
                        cell_data = PptxParser._extract_data_from_cell(table_cell)
                        tbl_row.append(cell_data)
                    parsed_table.append(tbl_row)
                prs_tables.append(parsed_table)
        prs_tables = self._parse_daily_update(prs_tables)
        self._data = prs_tables


class DailyUpdatePptxParser(PptxParser):
    @staticmethod
    def _find_key_by_value(table, row_index, col_index):
        """
        This function finds the title of the value (number) inside the table
        :param table - the table
        :param row_index - the index of the row
        :param col_index - the index of the column
        :return: the title of the value
        """
        try:
            # checks if there is a title on top of the value
            if table[row_index-1][col_index] and not table[row_index-1][col_index].replace(',','').isdigit():
                return table[row_index-1][col_index]
            # checks if there is a title on right of the value
            if table[row_index][col_index-1] and not table[row_index][col_index-1].replace(',','').isdigit():
                return table[row_index][col_index-1]
            raise ValueError(f'You have in the cell: {row_index},{col_index} a number without a title'
                             f'(titles are supposed to be on top or on the right of the number)')
        except IndexError:
            raise ValueError(f'You have in the cell: {row_index},{col_index} a number without a title'
                             f'(titles are supposed to be on top or on the right of the number)')

    @staticmethod
    def parse_file(tables):
        parsed_tables = list()
        for table in tables[DAILY_UPDATE_TABLE_BOTTOM_BUFFER:DAILY_UPDATE_TABLE_TOP_BUFFER]:
            table_values = list()
            table_keys = list()
            for row_index in range(len(table)):
                for col_index in range(len(table[0])):
                    if table[row_index][col_index].replace(',','').isdigit():
                        table_values.append(table[row_index][col_index])
                        table_keys.append(DailyUpdatePptxParser._find_key_by_value(table, row_index, col_index))
            parsed_tables.append([table_keys, table_values])
        return parsed_tables


class PdfParser(FileParser):
    """
    This class represents Pdf file parser.
    """
    def parse_file(self):
        try:
            pdf_tables = tabula.read_pdf(input_path=self.path,
                                         pages="all",
                                         stream=True,
                                         silent=True)

        except Exception:
            print(f"failed to read {os.path.basename(self.path)}")
            return

        for table_df in pdf_tables:
            # the loop will work only in the first case, where the headers are correct.
            self._parse_cities(table_df)
        self._parse_daily_update()

    def _parse_cities(self, table_df):
        """
        table type check: checks if the headers match the cities format.
        THIS IS AN EXAMPLE FOR A TABLE TYPE CHECK.
        :param table_df:
        :return: None
        """
        if CITIES_COLUMNS.issubset(set(table_df.columns.tolist())):
            print("cities parse")
            parser = CitiesPdfParser(self.path)
            self._data.append(parser.parse_file())
            self._output_dir = CITIES_OUTPUT_DIR

    def _parse_daily_update(self):
        """
        This function Checks if the pdf file is a file contains COVID-19 daily update.
        :return: None
        """
        if os.path.basename(self.path).startswith(DAILY_UPDATE_FILE_PREFIX):
            print("daily update parse")
            parser = DailyUpdatePdfParser(self.path)
            self._data = parser.parse_file()
            self._output_dir = DAILY_UPDATE_OUTPUT_DIR


class CitiesPdfParser(PdfParser):
    """
    This class represents a parser of a pdf file that contains COVID-19 data divided into cities.
    """

    def parse_file(self):
        translator = ParserTranslator(to_lang='en', from_lang='he')
        pdf_tables = tabula.read_pdf(input_path=self.path,
                                     pages="all",
                                     stream=True,
                                     silent=True)
        fixed_data = []
        headers = ["City_Name", "Population", "Infected"]
        fixed_data.append(headers)

        for data_df in pdf_tables:
            data_df = data_df.where(pd.notnull(data_df), None)
            list_data = data_df.values.tolist()

            for line in list_data:
                # Solves the problem that the data moved one column right in the middle of the file
                if line[1] is not None:
                    if line[1].replace(",", "").isdigit():
                        fixed_data.append([translator.translate_word(line[0]),
                                           line[1].replace(",", ""),
                                           str(line[2]).replace(",", "")])
                else:
                    fixed_data.append([translator.translate_word(line[0]),
                                       line[2].replace(",", ""),
                                       str(line[3]).replace(",", "")])
        return fixed_data


class DailyUpdatePdfParser(PdfParser):
    """
    This class represents a parser of a pdf file that contains COVID-19 daily update data.
    """
    def parse_file(self):
        fixed_pdf_tables = list()
        translated_pdf_tables = list()
        pdf_tables = tabula.read_pdf_with_template(
            input_path=self.path,
            template_path=DAILY_UPDATE_TEMPLATE_PATH
        )
        
        # first 3 tables of pdf_tables are from confirmed patients table (top)
        # last 3 tables of pdf_tables are from treatment table (bottom)
        
        pdf_table0 = DailyUpdatePdfParser._fix_critical_confirmed_table(pdf_tables[0])
        pdf_table3 = DailyUpdatePdfParser._fix_treatment_table(pdf_tables[3])
        
        # concatenating to form two tables - according to the pdf's format
        fixed_pdf_tables.append(pd.concat([pdf_tables[2], pdf_table0, pdf_tables[1]], axis=1))
        fixed_pdf_tables.append(pd.concat([pdf_tables[5], pdf_table3, pdf_tables[4]], axis=1))

        translator = ParserTranslator(to_lang='en', from_lang='he')
        for pdf_table in fixed_pdf_tables:
            temp_df = pd.DataFrame(data=[pdf_table.columns.tolist()] + pdf_table.values.tolist())
            temp_df = temp_df.applymap(lambda x: translator.translate_word(str(x)))
            translated_df = pd.DataFrame(columns=temp_df.values.tolist()[0], data=[temp_df.values.tolist()[1]])
            translated_pdf_tables.append(translated_df)

        return translated_pdf_tables
    
    @staticmethod
    def _fix_critical_confirmed_table(critical_confirmed_table):
        fixed_critical_confirmed_table = critical_confirmed_table.drop(critical_confirmed_table.index[0])
        return fixed_critical_confirmed_table.reset_index(drop=True)

    @staticmethod
    def _fix_treatment_table(treatment_table):
        fixed_treatment_table = pd.DataFrame(columns=["number", "state"],
            data=treatment_table.values.tolist() + [treatment_table.columns.tolist()])
        
        fixed_treatment_table = fixed_treatment_table.transpose().values.tolist()
        
        return pd.DataFrame(
            columns=fixed_treatment_table[1],
            data=[fixed_treatment_table[0]]
        )